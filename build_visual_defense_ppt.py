# -*- coding: utf-8 -*-
"""Generate VISUAL-RICH defense PPT for 点时成金 with diagrams, cards, timelines, metrics."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── palette ──
NAVY    = RGBColor(0x1A, 0x23, 0x3A)
NAVY2   = RGBColor(0x12, 0x18, 0x29)
GOLD    = RGBColor(0xC9, 0xA2, 0x27)
GOLD_L  = RGBColor(0xE5, 0xD4, 0x9E)
GOLD_D  = RGBColor(0xA6, 0x80, 0x00)
CREAM   = RGBColor(0xF9, 0xF7, 0xF1)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
TXT     = RGBColor(0x2B, 0x2B, 0x2B)
GRAY    = RGBColor(0x7A, 0x7A, 0x7A)
GRAY_L  = RGBColor(0xB0, 0xB0, 0xB0)
RED     = RGBColor(0xC0, 0x39, 0x2B)      # 涨
GREEN   = RGBColor(0x27, 0xAE, 0x60)      # 跌
BLUE    = RGBColor(0x2C, 0x80, 0xBA)
TEAL    = RGBColor(0x16, 0xA0, 0x85)
PURPLE  = RGBColor(0x8E, 0x44, 0xAD)
ORANGE  = RGBColor(0xE6, 0x7E, 0x22)

EA_FONT = "PingFang SC"
LAT_FONT = "Arial"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


# ═══════════════ helpers ═══════════════
def set_font(run, name=EA_FONT, size=None, bold=None, color=None):
    run.font.name = LAT_FONT
    if size is not None: run.font.size = Pt(size)
    if bold is not None: run.font.bold = bold
    if color is not None: run.font.color.rgb = color


def rect(slide, x, y, w, h, fill=None, line=None, lw=Pt(1), shape=MSO_SHAPE.RECTANGLE, radius=None):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    sp.shadow.inherit = False
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = line; sp.line.width = lw
    return sp


def txt(slide, x, y, w, h, text, sz=14, bold=False, color=TXT,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=EA_FONT, ls=1.15):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, ln in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = ls
        r = p.add_run(); r.text = ln
        set_font(r, font, sz, bold, color)
    return tb


def rounded_card(slide, x, y, w, h, fill, title, body="", title_sz=15, body_sz=11, title_c=WHITE, body_c=TXT):
    """Rounded rectangle card with title bar and body."""
    rect(slide, x, y, w, h, fill=fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(slide, x+Inches(0.18), y+Inches(0.12), w-Inches(0.36), Inches(0.38),
        title, sz=title_sz, bold=True, color=title_c)
    if body:
        txt(slide, x+Inches(0.18), y+Inches(0.48), w-Inches(0.36), h-Inches(0.58),
            body, sz=body_sz, color=body_c)


def icon_circle(slide, cx, cy, r, fill, label="", label_sz=10, label_c=WHITE):
    sp = rect(slide, cx-r, cy-r, r*2, r*2, fill=fill, shape=MSO_SHAPE.OVAL)
    if label:
        txt(slide, cx-r+Inches(0.05), cy-Inches(0.08), r*2-Inches(0.1), r*2,
            label, sz=label_sz, bold=True, color=label_c, align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE)


def arrow_down(slide, x, y, w, color=GOLD):
    """Simple down arrow via chevron."""
    sp = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, x, y, w, Inches(0.22))
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    sp.line.fill.background()


def arrow_right(slide, x, y, h, color=GOLD):
    sp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, Inches(0.35), h)
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    sp.line.fill.background()


def page_footer(slide, idx):
    txt(slide, Inches(0.45), Inches(7.06), Inches(7), Inches(0.32),
        "点时成金 · 黄金价格30分钟方向预测系统", sz=9, color=GRAY)
    txt(slide, Inches(12.3), Inches(7.06), Inches(0.85), Inches(0.32),
        f"{idx:02d}", sz=12, bold=True, color=GOLD_D, align=PP_ALIGN.RIGHT)


# ═══════════════ SLIDE 1 – COVER ═══════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, fill=NAVY)
rect(s, 0, Inches(2.3), SW, Inches(2.8), fill=NAVY2)
# gold lines
rect(s, Inches(1.0), Inches(2.22), Inches(11.33), Pt(2.5), fill=GOLD)
rect(s, Inches(1.0), Inches(5.05), Inches(11.33), Pt(2.5), fill=GOLD)
# decorative circles
icon_circle(s, Inches(11.8), Inches(0.8), Inches(0.55), GOLD_L, "", 0, WHITE)
icon_circle(s, Inches(1.0), Inches(6.2), Inches(0.35), GOLD_D, "", 0, WHITE)
txt(s, Inches(1.0), Inches(0.65), Inches(10), Inches(0.45),
    "PROJECT DEFENSE · 项目答辩汇报", sz=14, bold=True, color=GOLD)
txt(s, Inches(1.0), Inches(2.55), Inches(11), Inches(1.0),
    "点 时 成 金", sz=56, bold=True, color=WHITE)
txt(s, Inches(1.0), Inches(3.72), Inches(11), Inches(0.65),
    "黄金价格 30 分钟方向预测系统", sz=28, bold=True, color=GOLD)
txt(s, Inches(1.0), Inches(4.52), Inches(11), Inches(0.5),
    "基于 LLM 语义理解 × 时序量化模型的 AI 信号系统", sz=17, color=GOLD_L)
txt(s, Inches(1.0), Inches(5.25), Inches(11), Inches(0.45),
    "AI Signal System for XAU/USD 30-Min Direction Forecasting", sz=13, color=GRAY_L)
txt(s, Inches(1.0), Inches(6.75), Inches(11), Inches(0.35),
    "汇报人：项目团队          日期：2026 年 8 月", sz=13, color=GRAY)

# ═══════════════ SLIDE 2 – AGENDA (visual timeline) ═══════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, fill=CREAM)
rect(s, 0, 0, Inches(0.22), SH, fill=GOLD)
txt(s, Inches(0.8), Inches(0.45), Inches(10), Inches(0.8),
    "目录 · CONTENTS", sz=34, bold=True, color=NAVY)
rect(s, Inches(0.8), Inches(1.28), Inches(3.5), Pt(3), fill=GOLD)

agenda_items = [
    ("01", "项目背景与需求", "Background & Requirements"),
    ("02", "核心目标与创新点", "Core Objectives & Innovation"),
    ("03", "技术架构设计", "Technical Architecture"),
    ("04", "数据源与采集体系", "Data Sources & Collection"),
    ("05", "功能模块实现", "Module Implementation"),
    ("06", "关键问题与解决", "Challenges & Solutions"),
    ("07", "项目成果展示", "Results & Demo"),
    ("08", "总结与展望", "Summary & Outlook"),
]
for i, (num, cn, en) in enumerate(agenda_items):
    col = i // 4
    row = i % 4
    bx = Inches(0.8 + col * 6.2)
    by = Inches(1.85 + row * 1.3)
    # number circle
    icon_circle(s, bx + Inches(0.3), by + Inches(0.32), Inches(0.36), GOLD, num, 16, WHITE)
    # Chinese title
    txt(s, bx + Inches(0.85), by + Inches(0.12), Inches(5), Inches(0.42),
        cn, sz=19, bold=True, color=NAVY)
    # English subtitle
    txt(s, bx + Inches(0.85), by + Inches(0.54), Inches(5), Inches(0.32),
        en, sz=11, color=GRAY)
page_footer(s, 2)

# ═══════════════ SLIDE 3 – BACKGROUND (problem cards) ═══════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, fill=CREAM)
rect(s, 0, 0, SW, Inches(1.15), fill=NAVY)
rect(s, 0, Inches(1.15), SW, Pt(3), fill=GOLD)
txt(s, Inches(0.6), Inches(0.12), Inches(5), Inches(0.3),
    "BACKGROUND & REQUIREMENTS", sz=12, bold=True, color=GOLD)
txt(s,  Inches(0.6), Inches(0.40), Inches(12), Inches(0.65),
    "一、项目背景与需求", sz=28, bold=True, color=WHITE)

# Problem card row
problems = [
    (RED,    "定价复杂",    "黄金兼具商品·货币·避险·通胀对冲四重属性\n受美元指数/实际利率/地缘政治/央行购金等多因子交织驱动"),
    (ORANGE, "传统滞后",    "量化模型对突发事件反应滞后\n难以理解「鹰派降息」「避险升温」等复合语境"),
    (BLUE,   "极短时空白",  "黄金市场缺乏 30 分钟级事件驱动交易工具\n现有产品多为日频或更长周期"),
    (PURPLE, "语义鸿沟",    "新闻/讲话/公告中的隐含信息无法被数值化\nLLM 可桥接「自然语言 → 结构化因子」的缺口"),
]
for i, (clr, ttl, bd) in enumerate(problems):
    bx = Inches(0.45 + i * 3.2)
    by = Inches(1.55)
    rounded_card(s, bx, by, Inches(3.0), Inches(2.35), clr, ttl, bd, 15, 11)

# Solution box
rect(s, Inches(0.45), Inches(4.15), Inches(12.43), Inches(2.55),
     fill=WHITE, line=GOLD, lw=Pt(2), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, Inches(0.7), Inches(4.3), Inches(12), Inches(0.42),
    "💡 核心设想 · Core Idea", sz=17, bold=True, color=GOLD_D)
txt(s, Inches(0.7), Inches(4.78), Inches(11.8), Inches(1.8),
    "用大语言模型（LLM）实时解析新闻、官员讲话、地缘事件，提取情感分数与事件强度；\n"
    "结合多因子时序预测模型（LightGBM/XGBoost 集成），输出未来 30 分钟 XAU/USD 的\n"
    "上涨概率（0~1）、预估收益率、置信度及建议仓位——填补黄金市场极短时信号空白。",
    sz=15, color=TXT, ls=1.35)

# User icons at bottom
users = [
    ("👤", "短线交易员", "实时信号\n快速决策"),
    ("📊", "量化研究员", "回测分析\n准确率评估"),
    ("⚙️", "系统管理员", "监控运维\n数据质量"),
]
for i, (ic, nm, desc) in enumerate(users):
    ux = Inches(0.8 + i * 4.2)
    uy = Inches(6.82)
    icon_circle(s, ux + Inches(0.2), uy, Inches(0.28), NAVY, "", 1, WHITE)
    txt(s, ux + Inches(0.6), uy - Inches(0.04), Inches(3), Inches(0.28),
        f"{nm}  {desc}", sz=11, color=GRAY)
page_footer(s, 3)

# ═══════════════ SLIDE 4 – GOALS (metric cards) ═══════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, fill=CREAM)
rect(s, 0, 0, SW, Inches(1.15), fill=NAVY)
rect(s, 0, Inches(1.15), SW, Pt(3), fill=GOLD)
txt(s, Inches(0.6), Inches(0.12), Inches(5), Inches(0.3),
    "CORE OBJECTIVES", sz=12, bold=True, color=GOLD)
txt(s,  Inches(0.6), Inches(0.40), Inches(12), Inches(0.65),
    "二、核心目标与创新点", sz=28, bold=True, color=WHITE)

# Metric gauge cards
metrics = [
    (TEAL,   "> 60%",  "方向准确率",    "Direction Accuracy"),
    (BLUE,   "> 1.5",  "盈亏比",       "Profit/Loss Ratio"),
    (PURPLE, "> 1.0",  "夏普比率",     "Sharpe Ratio"),
    (ORANGE, "< 3s",   "信号延迟",     "Signal Latency"),
    (RED,    "< 15%",  "最大回撤",     "Max Drawdown"),
]
for i, (clr, val, cn, en) in enumerate(metrics):
    mx = Inches(0.45 + i * 2.56)
    my = Inches(1.5)
    rect(s, mx, my, Inches(2.38), Inches(1.95), fill=WHITE, line=clr, lw=Pt(2.5),
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    # colored top accent
    rect(s, mx, my, Inches(2.38), Inches(0.08), fill=clr)
    txt(s, mx, my + Inches(0.22), Inches(2.38), Inches(0.7),
        val, sz=32, bold=True, color=clr, align=PP_ALIGN.CENTER)
    txt(s, mx, my + Inches(0.88), Inches(2.38), Inches(0.38),
        cn, sz=15, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    txt(s, mx, my + Inches(1.26), Inches(2.38), Inches(0.32),
        en, sz=10, color=GRAY, align=PP_ALIGN.CENTER)

# Innovation cards
innovations = [
    ("① 多模态融合", "同时考虑货币政策、地缘风险、实际利率、\n美元指数对黄金的综合影响", NAVY),
    ("② 地缘实时量化", "利用 LLM 从新闻中提取风险强度，\n而非仅依赖滞后月频指数", TEAL),
    ("③ 极短时预测", "30 分钟级别方向预测，\n填补黄金市场事件驱动交易的空白", BLUE),
    ("④ LLM×量化互补", "大模型解决语义理解问题，\n时序模型解决价格预测问题", PURPLE),
]
for i, (ttl, bd, clr) in enumerate(innovations):
    ix = Inches(0.45 + i * 3.2)
    iy = Inches(3.7)
    rect(s, ix, iy, Inches(3.0), Inches(2.55), fill=clr, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, ix + Inches(0.18), iy + Inches(0.18), Inches(2.64), Inches(0.4),
        ttl, sz=15, bold=True, color=WHITE)
    txt(s, ix + Inches(0.18), iy + Inches(0.62), Inches(2.64), Inches(1.8),
        bd, sz=12, color=RGBColor(0xE8, 0xE8, 0xF0))

# bottom tagline
rect(s, Inches(0.45), Inches(6.48), Inches(12.43), Inches(0.42), fill=GOLD_L, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, Inches(0.45), Inches(6.50), Inches(12.43), Inches(0.4),
    "✦ 定位：开发 AI 信号系统，实现 30 分钟级别黄金价格方向预测",
    sz=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
page_footer(s, 4)

# ═══════════════ SLIDE 5 – ARCHITECTURE (layered diagram like screenshot) ═══════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, fill=CREAM)
rect(s, 0, 0, SW, Inches(1.15), fill=NAVY)
rect(s, 0, Inches(1.15), SW, Pt(3), fill=GOLD)
txt(s, Inches(0.6), Inches(0.12), Inches(5), Inches(0.3),
    "TECHNICAL ARCHITECTURE", sz=12, bold=True, color=GOLD)
txt(s,  Inches(0.6), Inches(0.40), Inches(12), Inches(0.65),
    "三、技术架构设计 — 五层流水线", sz=28, bold=True, color=WHITE)

# Layer boxes (bottom to top visually, but drawn top to bottom)
layers = [
    (Inches(1.4), Inches(5.55), RGBColor(0xFD, 0xED, 0xDA), ORANGE,
     "数据采集层 · 9 个独立爬虫包",
     ["因子爬虫 ×7  (DXY/TIPS/VIX/GPR/EPU/DFII/财经日历)",
      "XAUUSD 30m 爬虫  (新浪 hq.sinajs.cn, 30秒轮询)",
      "新闻爬虫 4站+LLM  (美联储/白宫/AP/CNN → DashScope qwen)",
      "统一模板结构: config/scraper/parser/storage/utils/main"]),
    (Inches(1.4), Inches(4.35), RGBColor(0xD5, 0xE8, 0xD4), TEAL,
     "数据访问层 · app/models",
     ["SQLAlchemy 2.0 ORM  ·  8 张业务表",
      "SQLite (WAL模式)  ·  data/gold_predictor.db",
      "方言自适应: PG ↔ SQLite 仅改一行 DATABASE_URL"]),
    (Inches(1.4), Inches(3.15), RGBColor(0xE8, 0xDA, 0xEF), PURPLE,
     "业务逻辑层 · app/core  (1558 行, 9 模块)",
     ["因子采集适配器  ·  新闻采集+去重(rapidfuzz>0.8)",
      "LLM 情感分析  ·  鹰鸽指数打分  ·  信号生成(6因子加权)",
      "Backtrader 回测  ·  数据质量检测  ·  历史预加载"]),
    (Inches(1.4), Inches(1.95), RGBColor(0xDF, 0xE2, 0xE8), BLUE,
     "接口层 · FastAPI",
     ["8 组 RESTful 路由  (/api/v1/signals|market|factors|news|...)",
      "统一 ApiResponse{code,message,data} 信封",
      "依赖注入(deps.py)  ·  CORS  ·  TestClient 13端点全200"]),
    (Inches(1.4), Inches(1.45), RGBColor(0xDE, 0xEB, 0xF7), NAVY,
     "表现层 · Presentation",
     ["Streamlit 仪表盘  app/dashboard  (5 大面板组件)",
      "静态原型  frontend/dashboard.html  (暖色财经终端风)",
      "DEMO_MODE 优雅降级: 后端不可用时读本地 demo JSON"]),
]

for lx, ly, bgc, brc, title, items in layers:
    bw = Inches(10.53)
    bh = Inches(0.78 + len(items) * 0.28)
    rect(s, lx, ly, bw, bh, fill=bgc, line=brc, lw=Pt(1.5), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, lx + Inches(0.15), ly + Inches(0.06), bw - Inches(0.3), Inches(0.34),
        title, sz=13, bold=True, color=brc)
    for j, it in enumerate(items):
        txt(s, lx + Inches(0.25), ly + Inches(0.4 + j * 0.28), bw - Inches(0.5), Inches(0.28),
            f"• {it}", sz=10.5, color=TXT)

# Side annotation: key design patterns
rect(s, Inches(12.0), Inches(1.5), Inches(1.1), Inches(4.8), fill=WHITE, line=GOLD, lw=Pt(1.5),
     shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, Inches(12.05), Inches(1.58), Inches(1.0), Inches(0.32),
    "核心设计", sz=10, bold=True, color=GOLD_D, align=PP_ALIGN.CENTER)
patterns = ["适配器模式", "幂等upsert", "三级降级", "防泄漏验证", "单向依赖"]
for j, pat in enumerate(patterns):
    icon_circle(s, Inches(12.47), Inches(2.05 + j * 0.85), Inches(0.22), GOLD, str(j+1), 9, WHITE)
    txt(s, Inches(12.03), Inches(2.2 + j * 0.85), Inches(1.0), Inches(0.55),
        pat, sz=9, color=NAVY, align=PP_ALIGN.CENTER)

page_footer(s, 5)

# ═══════════════ SLIDE 6 – DATA SOURCES (grid of cards) ═══════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, fill=CREAM)
rect(s, 0, 0, SW, Inches(1.15), fill=NAVY)
rect(s, 0, Inches(1.15), SW, Pt(3), fill=GOLD)
txt(s, Inches(0.6), Inches(0.12), Inches(5), Inches(0.3),
    "DATA SOURCES & COLLECTION", sz=12, bold=True, color=GOLD)
txt(s,  Inches(0.6), Inches(0.40), Inches(12), Inches(0.65),
    "四、数据源与采集体系", sz=28, bold=True, color=WHITE)

# Category headers + data source cards
data_grid = [
    # (category, color, [(name, source, freq, tech)])
    ("🌍 地缘政治类", BLUE, [
        ("GPR 地缘风险指数", "matteoiacoviello.com .xls", "月频", "requests+pandas"),
        ("VIX 恐慌指数", "CBOE 官网", "日频/实时", "Playwright sync"),
        ("WTI 原油价格", "EIA / ICE", "实时/日频", "API / 爬虫"),
        ("美元指数 DXY", "新浪财经", "日频/实时", "Playwright sync"),
    ]),
    ("🇺🇸 美国经济类", TEAL, [
        ("TIPS 实际利率", "FRED CSV", "日频", "subprocess curl"),
        ("美债 10Y-2Y 利差", "FRED API", "日频", "requests"),
        ("EPU 政策不确定", "policyuncertainty.com", "日频", "requests"),
        ("鹰鸽指数", "美联储/白宫讲话", "事件级", "LLM DashScope"),
    ]),
    ("📰 新闻与情感", PURPLE, [
        ("美联储官网", "federalreserve.gov", "事件级", "Playwright async"),
        ("白宫新闻", "whitehouse.gov", "事件级", "Playwright async"),
        ("AP / CNN", "apnews.com / CNN", "事件级", "Playwright async"),
        ("LLM 情感分析", "阿里云 DashScope qwen-turbo", "实时", "LangChain"),
    ]),
    ("📊 行情数据", ORANGE, [
        ("XAU/USD 30m K线", "新浪 hq.sinajs.cn", "30秒轮询", "urllib"),
        ("财经日历", "investing.com sslecal2", "事件级", "BS4+requests"),
        ("COMEX COT持仓", "CFTC 周报", "周频", "手动/API"),
        ("ETF 持仓变动", "WGC / 各发行商", "日频", "API / 爬虫"),
    ]),
]

start_y = Inches(1.4)
for ci, (cat_name, cat_clr, sources) in enumerate(data_grid):
    cy = start_y + ci * Inches(1.46)
    # category label
    rect(s, Inches(0.35), cy, Inches(2.0), Inches(1.35), fill=cat_clr, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, Inches(0.4), cy + Inches(0.45), Inches(1.9), Inches(0.5),
        cat_name, sz=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # source mini-cards
    for si, (nm, src, frq, tech) in enumerate(sources):
        sx = Inches(2.5) + si * Inches(2.68)
        rect(s, sx, cy, Inches(2.55), Inches(1.35), fill=WHITE, line=cat_clr, lw=Pt(1),
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        txt(s, sx + Inches(0.1), cy + Inches(0.08), Inches(2.35), Inches(0.34),
            nm, sz=11, bold=True, color=cat_clr)
        txt(s, sx + Inches(0.1), cy + Inches(0.42), Inches(2.35), Inches(0.26),
            src, sz=8.5, color=GRAY)
        txt(s, sx + Inches(0.1), cy + Inches(0.70), Inches(1.2), Inches(0.24),
            f"📅 {frq}", sz=9, color=TEAL)
        txt(s, sx + Inches(1.3), cy + Inches(0.70), Inches(1.2), Inches(0.24),
            f"⚙ {tech}", sz=8.5, color=GRAY, align=PP_ALIGN.RIGHT)

# Bottom note
rect(s, Inches(0.35), Inches(6.78), Inches(12.63), Inches(0.38), fill=GOLD_L, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, Inches(0.35), Inches(6.79), Inches(12.63), Inches(0.36),
    "统一落盘约定: OUTPUT_DIR=data/  →  X_latest.json (快照) + X_history.jsonl (追加历史)", sz=11, bold=True, color=NAVY,
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
page_footer(s, 6)

# ═══════════════ SLIDE 7 – MODULES (feature cards) ═══════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, fill=CREAM)
rect(s, 0, 0, SW, Inches(1.15), fill=NAVY)
rect(s, 0, Inches(1.15), SW, Pt(3), fill=GOLD)
txt(s, Inches(0.6), Inches(0.12), Inches(5), Inches(0.3),
    "MODULE IMPLEMENTATION", sz=12, bold=True, color=GOLD)
txt(s,  Inches(0.6), Inches(0.40), Inches(12), Inches(0.65),
    "五、主要功能模块实现", sz=28, bold=True, color=WHITE)

modules = [
    ("F01-F03", "数据采集", NAVY,
     "9 个爬虫覆盖行情/宏观/新闻/日历\n统一六件套模板 + 幂等落库\n快照(JSON) + 历史(JSONL) 双写"),
    ("F04-F06", "模型层·智能分析", TEAL,
     "LLM 新闻情感: −1(利空) ~ +1(利多)\n鹰鸽指数: 官员讲话 → 黄金影响因子\n地缘风险: 关键词提取 + 滑动窗口平滑"),
    ("F07", "时序预测模型", BLUE,
     "24 维特征: 技术10 + 情感10 + 市场4\nLightGBM/XGBoost 加权集成\nPurge Walk-Forward CV 防泄漏"),
    ("F08", "信号生成器", PURPLE,
     "每 5 分钟聚合 6 因子输出信号\n方向(涨/跌/观望) + 强度(0~100)\n仓位建议 + 止损止盈"),
    ("F09", "可视化仪表盘", ORANGE,
     "Streamlit 5 大面板: 行情/信号/资讯/因子/绩效\n前端静态原型(暖色财经终端)\nDEMO_MODE 降级可离线演示"),
    ("F10-F11", "回测与统计", RED,
     "Backtrader 模拟交易引擎\n7d/30d 滚动准确率统计\n权益曲线 + 最大回撤计算"),
]

for i, (code, title, clr, body) in enumerate(modules):
    col = i % 3
    row = i // 3
    mx = Inches(0.4 + col * 4.28)
    my = Inches(1.45 + row * 2.85)
    # card background
    rect(s, mx, my, Inches(4.08), Inches(2.65), fill=WHITE, line=clr, lw=Pt(2),
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    # code badge
    rect(s, mx + Inches(0.12), my + Inches(0.12), Inches(0.9), Inches(0.32), fill=clr,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, mx + Inches(0.12), my + Inches(0.12), Inches(0.9), Inches(0.32),
        code, sz=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # title
    txt(s, mx + Inches(1.1), my + Inches(0.12), Inches(2.85), Inches(0.36),
        title, sz=16, bold=True, color=clr)
    # divider
    rect(s, mx + Inches(0.12), my + Inches(0.52), Inches(3.84), Pt(1), fill=GRAY_L)
    # body
    txt(s, mx + Inches(0.18), my + Inches(0.65), Inches(3.72), Inches(1.9),
        body, sz=12, color=TXT, ls=1.3)

page_footer(s, 7)

# ═══════════════ SLIDE 8 – DATA FLOW PIPELINE ═══════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, fill=CREAM)
rect(s, 0, 0, SW, Inches(1.15), fill=NAVY)
rect(s, 0, Inches(1.15), SW, Pt(3), fill=GOLD)
txt(s, Inches(0.6), Inches(0.12), Inches(5), Inches(0.3),
    "DATA PIPELINE", sz=12, bold=True, color=GOLD)
txt(s,  Inches(0.6), Inches(0.40), Inches(12), Inches(0.65),
    "数据处理流水线 · End-to-End Flow", sz=28, bold=True, color=WHITE)

# Pipeline stages as connected blocks
pipeline = [
    (ORANGE, "① 采集", "9 个爬虫并行抓取\n去重(rapidfuzz)\n清洗标准化"),
    (TEAL,   "② 存储", "SQLite WAL 模式\n幂等 upsert\n8 张业务表"),
    (BLUE,   "③ 分析", "LLM 情感打分\n鹰鸽/地缘量化\n特征工程 24维"),
    (PURPLE, "④ 预测", "LightGBM/XGBoost\n集成学习\nWalk-Forward CV"),
    (GOLD_D, "⑤ 输出", "信号生成器\n方向/强度/仓位\n仪表盘展示"),
]

pw = Inches(2.3)
ph = Inches(1.85)
px_start = Inches(0.5)
py = Inches(1.6)
for i, (clr, stage, desc) in enumerate(pipeline):
    px = px_start + i * Inches(2.55)
    # main block
    rect(s, px, py, pw, ph, fill=clr, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, px, py + Inches(0.12), pw, Inches(0.4),
        stage, sz=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, px + Inches(0.1), py + Inches(0.55), pw - Inches(0.2), Inches(1.2),
        desc, sz=11, color=RGBColor(0xF0, 0xF0, 0xFF), align=PP_ALIGN.CENTER)
    # arrow between blocks
    if i < len(pipeline) - 1:
        arr_x = px + pw + Inches(0.02)
        arr_y = py + ph / 2 - Inches(0.15)
        sp = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, arr_x, arr_y, Inches(0.22), Inches(0.3))
        sp.fill.solid(); sp.fill.fore_color.rgb = GRAY_L
        sp.line.fill.background()

# Knowledge graph section below
rect(s, Inches(0.5), Inches(3.75), Inches(12.33), Inches(2.9), fill=WHITE, line=NAVY, lw=Pt(1.5),
     shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, Inches(0.7), Inches(3.85), Inches(12), Inches(0.38),
    "🔗 宏观传导知识图谱 · Causal Chain (部分)", sz=15, bold=True, color=NAVY)

chains = [
    ("实际利率 ↑", "→", "持有成本 ↑", "→", "金价 ↓", NAVY),
    ("美元指数 ↑", "→", "非美购金成本 ↑", "→", "金价 ↓", BLUE),
    ("通胀超预期", "→", "法币贬值预期", "→", "金价 ↑", RED),
    ("央行增持", "→", "需求增加信号", "→", "金价 ↑", TEAL),
    ("地缘冲突升级", "→", "避险情绪 ↑", "→", "金价 ↑", PURPLE),
]

for i, (a, ar, b, br, c, clr) in enumerate(chains):
    cy = Inches(4.35 + i * 0.44)
    # node A
    rect(s, Inches(0.75), cy, Inches(1.85), Inches(0.36), fill=clr, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, Inches(0.75), cy, Inches(1.85), Inches(0.36), a, sz=10, bold=True, color=WHITE,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # arrow text
    txt(s, Inches(2.7), cy, Inches(0.35), Inches(0.36), ar, sz=14, bold=True, color=GRAY,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # node B
    rect(s, Inches(3.1), cy, Inches(2.0), Inches(0.36), fill=clr, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, Inches(3.1), cy, Inches(2.0), Inches(0.36), b, sz=10, bold=True, color=WHITE,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # arrow text 2
    txt(s, Inches(5.2), cy, Inches(0.35), Inches(0.36), br, sz=14, bold=True, color=GRAY,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # result
    rect(s, Inches(5.6), cy, Inches(1.1), Inches(0.36), fill=clr, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, Inches(5.6), cy, Inches(1.1), Inches(0.36), c, sz=10, bold=True, color=WHITE,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# DB tables summary on right
txt(s, Inches(7.2), Inches(4.0), Inches(5.3), Inches(0.32),
    "📦 数据模型 · 8 张业务表", sz=13, bold=True, color=GOLD_D)
tables_text = (
    "news           — 新闻原文 (url 唯一去重)\n"
    "sentiment      — LLM 情感分析结果\n"
    "market_data    — 行情 (GC=F/DX-Y.NYB/^VIX/^IRX)\n"
    "factor_data    — 多因子 (DXY/TIPS/VIX/情感均值)\n"
    "signals        — 方向信号\n"
    "hawk_dove_events — 鹰鸽事件打分\n"
    "backtest_results — 回测结果\n"
    "economic_calendar — 财经日历事件")
txt(s, Inches(7.2), Inches(4.38), Inches(5.5), Inches(2.15),
    tables_text, sz=10, color=TXT, font="Menlo, monospace")

page_footer(s, 8)

# ═══════════════ SLIDE 9 – CHALLENGES (problem-solution pairs) ═══════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, fill=CREAM)
rect(s, 0, 0, SW, Inches(1.15), fill=NAVY)
rect(s, 0, Inches(1.15), SW, Pt(3), fill=GOLD)
txt(s, Inches(0.6), Inches(0.12), Inches(5), Inches(0.3),
    "CHALLENGES & SOLUTIONS", sz=12, bold=True, color=GOLD)
txt(s,  Inches(0.6), Inches(0.40), Inches(12), Inches(0.65),
    "六、关键问题与解决思路", sz=28, bold=True, color=WHITE)

challenges = [
    ("模块名冲突", "9个爬虫用裸名模块互相污染 sys.modules",
     "隔离动态加载", "在隔离 sys.path 下逐个加载并主动清除同名残留模块"),
    ("数据源异构", "FRED/新浪/CBOE 返回格式不统一",
     "方言自适应", "JSONB→JSON / Uuid→CHAR(32) 可移植适配，改一行切换 PG↔SQLite"),
    ("数据泄漏", "随机切分高估模型性能",
     "时序防漏", "Purge/Embargo + Purged Walk-Forward CV，超参仅在开发集内选"),
    ("LLM 可靠性", "全文重算贵、结构化输出易失败",
     "三级降级", "去重跳过 + Semaphore(3)限流 + 最终降级为'中性+低置信'"),
    ("后端分支分叉", "工作区与团队主干各缺模块",
     "合并整合", "合并 FastAPI 后端(app/api/ 8组路由)，消除分叉，TestClient 全200"),
    ("联调缺陷", "DB不同步/颜色反转/回测越界",
     "逐项修复", "重建种子库 + inv()颜色转换 + 重写回测引擎，jsdom 0错误"),
]

for i, (prob, prob_desc, sol, sol_desc) in enumerate(challenges):
    col = i % 2
    row = i // 2
    cx = Inches(0.4 + col * 6.5)
    cy = Inches(1.45 + row * 1.85)
    # problem (left/red-ish)
    rect(s, cx, cy, Inches(3.0), Inches(1.65), fill=RGBColor(0xFD, 0xED, 0xED), line=RED, lw=Pt(1.5),
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, cx + Inches(0.12), cy + Inches(0.08), Inches(2.76), Inches(0.32),
        f"⚠ {prob}", sz=13, bold=True, color=RED)
    txt(s, cx + Inches(0.12), cy + Inches(0.44), Inches(2.76), Inches(1.1),
        prob_desc, sz=10.5, color=TXT)
    # arrow
    ax = cx + Inches(3.05)
    ay = cy + Inches(0.6)
    sp = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, ax, ay, Inches(0.32), Inches(0.35))
    sp.fill.solid(); sp.fill.fore_color.rgb = GREEN
    sp.line.fill.background()
    # solution (right/green-ish)
    sx = cx + Inches(3.45)
    rect(s, sx, cy, Inches(2.9), Inches(1.65), fill=RGBColor(0xE2, 0xF0, 0xE5), line=GREEN, lw=Pt(1.5),
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, sx + Inches(0.12), cy + Inches(0.08), Inches(2.66), Inches(0.32),
        f"✓ {sol}", sz=13, bold=True, color=GREEN)
    txt(s, sx + Inches(0.12), cy + Inches(0.44), Inches(2.66), Inches(1.1),
        sol_desc, sz=10.5, color=TXT)

page_footer(s, 9)

# ═══════════════ SLIDE 10 – RESULTS (dashboard mockup + metrics) ═══════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, fill=CREAM)
rect(s, 0, 0, SW, Inches(1.15), fill=NAVY)
rect(s, 0, Inches(1.15), SW, Pt(3), fill=GOLD)
txt(s, Inches(0.6), Inches(0.12), Inches(5), Inches(0.3),
    "RESULTS & DEMO", sz=12, bold=True, color=GOLD)
txt(s,  Inches(0.6), Inches(0.40), Inches(12), Inches(0.65),
    "七、项目成果与演示", sz=28, bold=True, color=WHITE)

# Left side: Dashboard modules mockup
rect(s, Inches(0.4), Inches(1.4), Inches(6.3), Inches(5.4), fill=WHITE, line=GRAY_L, lw=Pt(1),
     shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, Inches(0.55), Inches(1.5), Inches(6), Inches(0.36),
    "🖥 仪表盘 5 大模块 (可现场演示)", sz=14, bold=True, color=NAVY)

dash_modules = [
    ("① 实时行情走势", "XAU/USD 价格曲线 · 涨跌幅 · 当前价", BLUE),
    ("② 信号方向归因", "看涨/看跌/观望 · 概率 · 强度 · 多空评分 · 因子贡献", TEAL),
    ("③ 资讯情感标签", "新闻列表 · 情感标签(红涨绿跌) · 关键句 · 置信度", PURPLE),
    ("④ 多因子实时值", "DXY/TIPS/VIX/情感均值 · trend_color · 变动量", ORANGE),
    ("⑤ 绩效回测曲线", "策略 vs 基准权益曲线 · 准确率 · 最大回撤 · 夏普比率", GREEN),
]
for i, (nm, desc, clr) in enumerate(dash_modules):
    dy = Inches(1.95 + i * 0.94)
    rect(s, Inches(0.55), dy, Inches(6.0), Inches(0.84), fill=RGBColor(0xFA, 0xFA, 0xFA), line=clr, lw=Pt(1),
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    # colored left bar
    rect(s, Inches(0.55), dy, Inches(0.06), Inches(0.84), fill=clr)
    txt(s, Inches(0.72), dy + Inches(0.06), Inches(5.7), Inches(0.32),
        nm, sz=12, bold=True, color=clr)
    txt(s, Inches(0.72), dy + Inches(0.4), Inches(5.7), Inches(0.38),
        desc, sz=10, color=GRAY)

# Right side: Key numbers / stats
rect(s, Inches(6.9), Inches(1.4), Inches(6.0), Inches(2.6), fill=NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, Inches(7.05), Inches(1.5), Inches(5.7), Inches(0.36),
    "📊 工程成果指标", sz=14, bold=True, color=GOLD)

eng_stats = [
    ("13", "API 端点", "TestClient 全部返回 200"),
    ("34", "测试用例", "pytest 全部通过"),
    ("11", "接口联调", "10 GET + 1 POST 全打通"),
    ("0", "运行时错误", "jsdom 全视图加载验证"),
]
for i, (val, label, note) in enumerate(eng_stats):
    ex = Inches(7.1 + (i % 2) * 2.85)
    ey = Inches(1.95 + (i // 2) * 0.98)
    txt(s, ex, ey, Inches(1.0), Inches(0.5), val, sz=28, bold=True, color=GOLD)
    txt(s, ex + Inches(0.9), ey + Inches(0.04), Inches(1.8), Inches(0.32),
        label, sz=12, bold=True, color=WHITE)
    txt(s, ex, ey + Inches(0.48), Inches(2.7), Inches(0.28),
        note, sz=9, color=GRAY_L)

# Model training results
rect(s, Inches(6.9), Inches(4.15), Inches(6.0), Inches(2.65), fill=WHITE, line=TEAL, lw=Pt(1.5),
     shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, Inches(7.05), Inches(4.25), Inches(5.7), Inches(0.36),
    "🧪 模型训练流水线", sz=14, bold=True, color=TEAL)

model_stats = [
    ("124,272", "根 30m K 线", "真实锚定 (2016~2026 日频拆出)"),
    ("120", "条真实情感", "LLM 文件 + SQLite 爬取新闻"),
    ("6,790", "日 GPR/EPU", "真实地缘/政策不确定性指数"),
    ("0.9156", "WF AUC", "Purged Walk-Forward 3 折 ±0.0065"),
]
for i, (val, label, note) in enumerate(model_stats):
    ex = Inches(7.1 + (i % 2) * 2.85)
    ey = Inches(4.7 + (i // 2) * 0.98)
    txt(s, ex, ey, Inches(1.2), Inches(0.42), val, sz=20, bold=True, color=TEAL)
    txt(s, ex + Inches(1.15), ey + Inches(0.04), Inches(1.6), Inches(0.28),
        label, sz=11, bold=True, color=NAVY)
    txt(s, ex, ey + Inches(0.44), Inches(2.7), Inches(0.28),
        note, sz=8.5, color=GRAY)

# Honesty disclaimer
rect(s, Inches(0.4), Inches(6.92), Inches(12.53), Inches(0.42), fill=RGBColor(0xFE, 0xF3, 0xCD),
     line=ORANGE, lw=Pt(1), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, Inches(0.5), Inches(6.93), Inches(12.33), Inches(0.4),
    "⚠️ 如实说明：当前 0.92 AUC 为「真实锚定数据的日内构造假象」，非真实 alpha。需接入真实 tick 级 30m K 线后重训方可作为效果依据。",
    sz=11, bold=False, color=ORANGE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

page_footer(s, 10)

# ═══════════════ SLIDE 11 – TIMELINE / ROADMAP ═══════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, fill=CREAM)
rect(s, 0, 0, SW, Inches(1.15), fill=NAVY)
rect(s, 0, Inches(1.15), SW, Pt(3), fill=GOLD)
txt(s, Inches(0.6), Inches(0.12), Inches(5), Inches(0.3),
    "TIMELINE & ROADMAP", sz=12, bold=True, color=GOLD)
txt(s,  Inches(0.6), Inches(0.40), Inches(12), Inches(0.65),
    "八、总结与展望", sz=28, bold=True, color=WHITE)

# Timeline - completed phase
rect(s, Inches(0.5), Inches(1.45), Inches(12.33), Inches(2.35), fill=WHITE, line=TEAL, lw=Pt(1.5),
     shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, Inches(0.65), Inches(1.53), Inches(4), Inches(0.34),
    "✅ 已完成 · Completed", sz=14, bold=True, color=TEAL)

done_items = [
    ("SRS 需求规格", "11项功能(F01~F11) + 四类非功能需求"),
    ("技术选型", "Python/FastAPI/SQLite/LightGBM/DashScope/Streamlit"),
    ("9 个爬虫", "行情/宏观/新闻全覆盖，统一模板+幂等落库"),
    ("FastAPI 后端", "8组路由 13端点 TestClient 全200"),
    ("前后端联调", "11接口打通 jsdom 0错误 颜色语义正确"),
    ("模型流水线", "12.4万bar+120情感+WF-CV+AutoML 跑通"),
]
for i, (nm, desc) in enumerate(done_items):
    dx = Inches(0.7 + (i % 3) * 4.05)
    dy = Inches(1.95 + (i // 3) * 0.9)
    icon_circle(s, dx + Inches(0.12), dy + Inches(0.22), Inches(0.18), TEAL, "✓", 9, WHITE)
    txt(s, dx + Inches(0.38), dy + Inches(0.02), Inches(3.5), Inches(0.3),
        nm, sz=11, bold=True, color=NAVY)
    txt(s, dx + Inches(0.38), dy + Inches(0.32), Inches(3.5), Inches(0.3),
        desc, sz=9, color=GRAY)

# Timeline - future roadmap
rect(s, Inches(0.5), Inches(3.95), Inches(12.33), Inches(2.9), fill=WHITE, line=ORANGE, lw=Pt(1.5),
     shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, Inches(0.65), Inches(4.03), Inches(4), Inches(0.34),
    "🔮 展望 · Roadmap", sz=14, bold=True, color=ORANGE)

roadmap = [
    ("P0", "真实 30m 数据", "接入 Yahoo/Stooq tick 级 K 线，替换锚定模拟路径", "获得可信 AUC"),
    ("P0", "新闻密度提升", "实时爬取 + 时间戳精确到分钟级", "验证情感增量"),
    ("P1", "SHAP 特征归因", "排列重要性 + 跨折稳定性检验", "定位真实信号"),
    ("P1", "WebSocket 推送", "替换 5 分钟轮询，降低延迟", "实时性 <1s"),
    ("P2", "鉴权 + CORS 收紧", "JWT / API Key 认证", "安全上线"),
    ("P2", "Docker 化部署", "单容器一键启动，生产就绪", "可交付运行"),
]
for i, (pri, nm, desc, goal) in enumerate(roadmap):
    rx = Inches(0.7 + (i % 3) * 4.05)
    ry = Inches(4.45 + (i // 3) * 1.12)
    # priority badge
    badge_clr = RED if pri == "P0" else (ORANGE if pri == "P1" else GRAY)
    rect(s, rx, ry, Inches(0.45), Inches(0.26), fill=badge_clr, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, rx, ry, Inches(0.45), Inches(0.26), pri, sz=8, bold=True, color=WHITE,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, rx + Inches(0.52), ry, Inches(3.3), Inches(0.3),
        nm, sz=11.5, bold=True, color=NAVY)
    txt(s, rx + Inches(0.52), ry + Inches(0.3), Inches(3.3), Inches(0.3),
        desc, sz=9, color=GRAY)
    txt(s, rx + Inches(0.52), ry + Inches(0.58), Inches(3.3), Inches(0.28),
        f"→ {goal}", sz=9, color=TEAL)

# Bottom takeaway
rect(s, Inches(0.5), Inches(7.0), Inches(12.33), Inches(0.35), fill=GOLD, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, Inches(0.5), Inches(7.0), Inches(12.33), Inches(0.35),
    "把「工程确定性」留给今天，把「真实 alpha」交给下一阶段的数据与训练。",
    sz=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

page_footer(s, 11)

# ═══════════════ SLIDE 12 – THANKS / Q&A ═══════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, fill=NAVY)
rect(s, Inches(1.2), Inches(2.8), Inches(11), Pt(2.5), fill=GOLD)
rect(s, Inches(1.2), Inches(4.5), Inches(11), Pt(2.5), fill=GOLD)
txt(s, Inches(1.2), Inches(3.0), Inches(11), Inches(1.0),
    "感谢聆听  ·  敬请指正", sz=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, Inches(1.2), Inches(4.65), Inches(11), Inches(0.6),
    "Q & A", sz=24, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
# contact info area
txt(s, Inches(1.2), Inches(5.5), Inches(11), Inches(0.4),
    "点时成金 · 黄金价格 30 分钟方向预测系统", sz=14, color=GOLD_L, align=PP_ALIGN.CENTER)
txt(s, Inches(1.2), Inches(5.95), Inches(11), Inches(0.35),
    "AI Signal System for XAU/USD 30-Min Direction Forecasting", sz=11, color=GRAY,
    align=PP_ALIGN.CENTER)

# Save
out_path = "/Users/echo/Desktop/forecasts for gold prices/点时成金_项目答辩_视觉版.pptx"
prs.save(out_path)
print(f"saved: {out_path}")
print(f"total slides: {len(prs.slides._sldIdLst)}")
