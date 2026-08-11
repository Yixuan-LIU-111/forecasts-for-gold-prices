# -*- coding: utf-8 -*-
"""Generate the project-defense PPT for 点时成金 (gold 30-min direction prediction)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------------- palette ----------------
DARK   = RGBColor(0x1F, 0x2A, 0x44)   # navy
DARK2  = RGBColor(0x14, 0x1C, 0x30)
GOLD   = RGBColor(0xC9, 0xA2, 0x27)
GOLD_D = RGBColor(0xB8, 0x86, 0x0B)
LIGHT  = RGBColor(0xF7, 0xF4, 0xEC)   # cream
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
TEXT   = RGBColor(0x2B, 0x2B, 0x2B)
GRAY   = RGBColor(0x6B, 0x6B, 0x6B)
SOFT   = RGBColor(0xEF, 0xEA, 0xDD)

EA_FONT = "PingFang SC"
LAT_FONT = "Arial"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def set_font(run, name=EA_FONT, size=None, bold=None, color=None, italic=None):
    run.font.name = LAT_FONT
    if size is not None:
        run.font.size = Pt(size)
    if bold is not None:
        run.font.bold = bold
    if italic is not None:
        run.font.italic = italic
    if color is not None:
        run.font.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", name)


def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=None, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = line_w or Pt(1)
    return sp


def add_text(slide, x, y, w, h, text, size=18, bold=False, color=TEXT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=EA_FONT,
             line_spacing=1.15, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(2); tf.margin_right = Pt(2)
    tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = ln
        set_font(r, name=font, size=size, bold=bold, color=color, italic=italic)
    return tb


def bg(slide, color=LIGHT):
    add_rect(slide, 0, 0, SW, SH, fill=color)


def footer(slide, idx):
    add_text(slide, Inches(0.4), Inches(7.04), Inches(7), Inches(0.35),
             "点时成金 · 黄金价格 30 分钟方向预测系统  |  项目答辩", size=9, color=GRAY)
    add_text(slide, Inches(12.2), Inches(7.04), Inches(0.9), Inches(0.35),
             f"{idx:02d}", size=11, bold=True, color=GOLD_D, align=PP_ALIGN.RIGHT)


def content_slide(title, bullets, idx, kicker=None):
    """bullets: list of (level, text)  level 0 = section bullet, 1 = sub."""
    s = prs.slides.add_slide(BLANK)
    bg(s, LIGHT)
    # left accent bar
    add_rect(s, 0, 0, Inches(0.22), SH, fill=GOLD)
    # title band
    add_rect(s, Inches(0.22), 0, Inches(13.111), Inches(1.18), fill=DARK)
    add_rect(s, Inches(0.22), Inches(1.18), Inches(13.111), Pt(3), fill=GOLD)
    if kicker:
        add_text(s, Inches(0.7), Inches(0.16), Inches(12), Inches(0.3),
                 kicker, size=12, bold=True, color=GOLD)
    add_text(s, Inches(0.7), Inches(0.42), Inches(12.2), Inches(0.7),
             title, size=26, bold=True, color=WHITE)
    # bullets
    tb = s.shapes.add_textbox(Inches(0.85), Inches(1.55), Inches(11.7), Inches(5.2))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (lvl, txt) in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if lvl == 0:
            p.space_before = Pt(10); p.space_after = Pt(2)
            p.line_spacing = 1.12
            r0 = p.add_run(); r0.text = "▪  "
            set_font(r0, size=16, bold=True, color=GOLD_D)
            r = p.add_run(); r.text = txt
            set_font(r, size=16.5, bold=True, color=DARK)
        else:
            p.space_before = Pt(2); p.space_after = Pt(6)
            p.line_spacing = 1.1
            p.level = 1
            r0 = p.add_run(); r0.text = "–  "
            set_font(r0, size=14, bold=False, color=GOLD)
            r = p.add_run(); r.text = txt
            set_font(r, size=14, color=TEXT)
    footer(s, idx)
    return s


# =========================================================
# Slide 1 — Cover
# =========================================================
s = prs.slides.add_slide(BLANK)
bg(s, DARK)
add_rect(s, 0, Inches(2.55), SW, Inches(2.4), fill=DARK2)
# gold rules
add_rect(s, Inches(1.2), Inches(2.42), Inches(11), Pt(2), fill=GOLD)
add_rect(s, Inches(1.2), Inches(4.98), Inches(11), Pt(2), fill=GOLD)
add_text(s, Inches(1.2), Inches(0.7), Inches(11), Inches(0.5),
         "PROJECT DEFENSE · 项目答辩", size=15, bold=True, color=GOLD)
add_text(s, Inches(1.2), Inches(2.7), Inches(11), Inches(1.1),
         "点时成金", size=52, bold=True, color=WHITE)
add_text(s, Inches(1.2), Inches(3.75), Inches(11), Inches(0.7),
         "黄金价格 30 分钟方向预测系统", size=30, bold=True, color=GOLD)
add_text(s, Inches(1.2), Inches(5.2), Inches(11), Inches(0.6),
         "基于 LLM 语义理解 + 时序量化的 AI 信号系统", size=17, color=LIGHT)
add_text(s, Inches(1.2), Inches(6.7), Inches(11), Inches(0.4),
         "汇报人：项目团队    日期：2026.08", size=13, color=GRAY)

# =========================================================
# Slide 2 — Agenda
# =========================================================
s = prs.slides.add_slide(BLANK)
bg(s, LIGHT)
add_rect(s, 0, 0, Inches(0.22), SH, fill=GOLD)
add_text(s, Inches(0.9), Inches(0.6), Inches(11), Inches(0.9),
         "目录 · CONTENTS", size=32, bold=True, color=DARK)
add_rect(s, Inches(0.9), Inches(1.5), Inches(3.2), Pt(3), fill=GOLD)

agenda = [
    ("01", "项目背景与需求"),
    ("02", "核心目标"),
    ("03", "技术架构与方案设计"),
    ("04", "主要功能模块及实现"),
    ("05", "开发关键问题与解决思路"),
    ("06", "项目成果与演示"),
    ("07", "总结与展望"),
]
col_x = [Inches(0.9), Inches(7.0)]
for i, (num, name) in enumerate(agenda):
    cx = col_x[i // 4]
    cy = Inches(2.0 + (i % 4) * 1.15)
    add_text(s, cx, cy, Inches(0.9), Inches(0.8), num, size=30, bold=True, color=GOLD)
    add_text(s, cx + Inches(1.0), cy + Inches(0.12), Inches(4.5), Inches(0.7),
             name, size=20, bold=True, color=TEXT, anchor=MSO_ANCHOR.MIDDLE)
footer(s, 2)

# =========================================================
# Content slides
# =========================================================
content_slide(
    "一、项目背景与需求",
    [
        (0, "黄金多重属性，定价高度复杂：兼具商品 / 货币 / 避险资产 / 通胀对冲属性，价格受美元指数、实际利率、地缘政治、央行购金、通胀预期等多因子交织驱动。"),
        (0, "传统量化模型的痛点：对突发事件反应滞后，难以理解“鹰派降息”“避险情绪升温”等复合语境。"),
        (0, "核心设想：用大语言模型（LLM）实时解析新闻 / 官员讲话 / 地缘事件，结合时序预测模型，输出未来 30 分钟 XAU/USD 上涨概率与预估幅度。"),
        (0, "三类目标用户：短线交易员（实时信号、快速决策）、量化研究员（回测与准确率分析）、系统管理员（监控与数据质量）。"),
        (0, "需求规格（SRS）：11 项功能（F01–F11）+ 多源数据分类（地缘 / 美国宏观 / 佐证表象 / 长期因素）+ 性能·准确·可用·安全四类非功能需求。"),
    ],
    idx=3, kicker="BACKGROUND & REQUIREMENTS"
)

content_slide(
    "二、核心目标",
    [
        (0, "核心定位：开发 AI 信号系统，实现 30 分钟级别黄金价格方向预测。"),
        (0, "量化验收指标：方向准确率 > 60% ｜ 盈亏比 > 1.5 ｜ 夏普比率 > 1.0 ｜ 信号延迟 < 3 秒 ｜ 最大回撤 < 15%。"),
        (0, "工程性能指标：信号延迟 < 3s、单条新闻情感分析 < 500ms、仪表盘加载 < 3s、系统可用性 > 99%（交易时段）。"),
        (0, "四大创新点："),
        (1, "多模态利多 / 利空融合：综合货币政策、地缘风险、实际利率、美元指数。"),
        (1, "地缘政治实时量化：用 LLM 从新闻提取风险强度，而非依赖滞后指数。"),
        (1, "极短时预测（30 分钟）：填补黄金市场事件驱动交易的空白。"),
        (1, "LLM + 传统量化互补：大模型解决语义理解，时序模型解决价格预测。"),
    ],
    idx=4, kicker="CORE OBJECTIVES"
)

content_slide(
    "三、技术架构与方案设计",
    [
        (0, "总体架构：数据层 → 模型层 → 应用层，三层单向依赖、边界清晰，无循环依赖。"),
        (0, "技术栈：Python 3.13 · FastAPI · SQLAlchemy/SQLite(WAL) · LightGBM/XGBoost/sklearn · LangChain + 阿里云 DashScope(qwen) · Streamlit · Backtrader · Docker。"),
        (0, "宏观传导知识图谱（因果链）：实际利率↑→金价↓、美元↑→金价↓、通胀超预期→金价↑、央行增持→金价↑、地缘冲突→避险→金价↑。"),
        (0, "关键设计决策："),
        (1, "适配器模式：隔离动态加载，解决 9 个爬虫同名模块冲突。"),
        (1, "方言自适应数据访问：PG ↔ SQLite 仅改一行 DATABASE_URL 即可切换。"),
        (1, "幂等 upsert + 真实唯一约束：防重复写入、保证可移植。"),
        (1, "三处优雅降级（仪表盘 / LLM / 爬虫）：保证流水线不中断。"),
        (1, "Purge/Embargo + Walk-Forward 时序验证：从源头防止数据泄漏。"),
    ],
    idx=5, kicker="ARCHITECTURE & DESIGN"
)

content_slide(
    "四、主要功能模块及实现",
    [
        (0, "数据采集层（9 个独立爬虫）：XAU/USD 30m、美元指数 DXY、TIPS 实际利率、VIX 恐慌指数、GPR 地缘风险、EPU 政策不确定、财经日历、新闻(+LLM 情感)、页面金价；统一“快照 + 历史”落盘约定。"),
        (0, "模型层：LLM 新闻情感分析（利多/利空 −1~+1）、鹰鸽指数、地缘风险量化、30 分钟多因子时序预测（LightGBM/XGBoost 加权集成）。"),
        (0, "应用层：信号生成器（每 5 分钟，方向/强度/仓位/止损止盈）、可视化仪表盘（5 大模块）、Backtrader 模拟回测。"),
        (0, "数据模型：8 张业务表 —— news / sentiment / market_data / factor_data / signals / hawk_dove_events / backtest_results / economic_calendar。"),
        (0, "后端整合：FastAPI 8 组路由（+deps）、app/main.py、统一 ApiResponse 信封，13 个端点经 TestClient 全部返回 200。"),
    ],
    idx=6, kicker="MODULES & IMPLEMENTATION"
)

content_slide(
    "五、开发过程中的关键问题与解决思路",
    [
        (0, "模块名冲突：9 个爬虫用裸名模块互相污染 sys.modules → 在隔离 sys.path 下逐个动态加载并主动清除同名残留模块。"),
        (0, "数据源异构：FRED / 新浪 / CBOE 等返回格式不一 → 方言自适应 upsert，JSONB→JSON、Uuid→CHAR(32) 等可移植适配。"),
        (0, "模型数据泄漏风险：随机切分会高估性能 → 时序 70/15/15 + Purge/Embargo + Purged Walk-Forward CV，超参仅在开发集内选择。"),
        (0, "LLM 成本与可靠性：全文重算昂贵、结构化输出易失败 → 按 (url,版本,模式) 去重跳过 + asyncio.Semaphore(3) 限流 + 三级降级到“中性+低置信”。"),
        (0, "后端分支分叉：工作区与团队主干各缺模块 → 合并 FastAPI 后端，恢复 API 层、消除分叉。"),
        (0, "联调缺陷：DB/ORM 不同步、涨跌颜色语义反转（国际 vs 国内）、回测回撤越界 → 重建种子库、inv() 转换、重写回测引擎。"),
    ],
    idx=7, kicker="KEY CHALLENGES & SOLUTIONS"
)

content_slide(
    "六、项目成果与演示",
    [
        (0, "工程成果：FastAPI 后端 13 个端点（/health + 12 个 /api/v1/*）TestClient 全 200；pytest 34 项全部通过。"),
        (0, "前后端联调：11 个接口（10 GET + 1 POST）全部打通，jsdom 全视图加载 0 运行时错误；字段映射、加载/空/错误态均验证通过。"),
        (0, "可视化仪表盘（5 大模块可演示）：① 实时行情走势　② 信号方向与多因子归因　③ 资讯情感标签　④ 多因子实时数值　⑤ 绩效回测曲线。"),
        (0, "模型训练流水线已跑通：真实锚定 30m 共 12.4 万根 bar、120 条真实新闻情感、6,790 日真实 GPR/EPU 压力；Purged Walk-Forward CV + AutoML 已落地。"),
        (0, "诚实的局限（答辩如实陈述）：当前 0.92 的 AUC 是“真实锚定数据的日内构造假象”，非真实 alpha；需接入真实 tick 级 30m K 线后重训方可作为效果依据。"),
    ],
    idx=8, kicker="RESULTS & DEMO"
)

content_slide(
    "七、总结与展望",
    [
        (0, "总结：项目已构建完整“多源采集 → 因子工程 → 方向信号 → 可视化”量化流水线，采集 / 模型 / 应用 / 回测工程链路端到端打通，架构清晰、可扩展、可演示。"),
        (0, "当前短板：真实 30m 行情与宏观因子覆盖不足，模型真实预测力待验证，缺写接口与鉴权。"),
        (0, "展望 ①：接入真实 tick 级 30m K 线，在真实数据上重训，并做 SHAP / 排列重要性定位真实信号。"),
        (0, "展望 ②：提升真实新闻事件密度与时刻精度，验证情感特征是否真有增量。"),
        (0, "展望 ③：WebSocket/SSE 实时推送、补齐写接口与鉴权、Docker 化部署上线。"),
        (0, "一句话收尾：把“工程确定性”留给今天，把“真实 alpha”交给下一阶段的数据与训练。"),
    ],
    idx=9, kicker="SUMMARY & OUTLOOK"
)

# =========================================================
# Slide 10 — Thanks / Q&A
# =========================================================
s = prs.slides.add_slide(BLANK)
bg(s, DARK)
add_rect(s, Inches(1.2), Inches(3.0), Inches(11), Pt(2), fill=GOLD)
add_rect(s, Inches(1.2), Inches(4.5), Inches(11), Pt(2), fill=GOLD)
add_text(s, Inches(1.2), Inches(3.2), Inches(11), Inches(1.0),
         "感谢聆听 · 敬请指正", size=40, bold=True, color=WHITE)
add_text(s, Inches(1.2), Inches(4.65), Inches(11), Inches(0.6),
         "Q & A", size=22, bold=True, color=GOLD)

out = "/Users/echo/Desktop/forecasts for gold prices/点时成金_项目答辩.pptx"
prs.save(out)
print("saved:", out, "slides:", len(prs.slides._sldIdLst))
